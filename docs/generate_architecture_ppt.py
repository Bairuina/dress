from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BLUE = RGBColor(37, 99, 235)
DARK = RGBColor(17, 24, 39)
GRAY = RGBColor(75, 85, 99)
LIGHT = RGBColor(243, 244, 246)
GREEN = RGBColor(22, 163, 74)
ORANGE = RGBColor(245, 158, 11)
RED = RGBColor(220, 38, 38)
PURPLE = RGBColor(168, 85, 247)
WHITE = RGBColor(255, 255, 255)
FONT_NAME = "Microsoft YaHei"


def zh(text: str) -> str:
    return text.encode("ascii").decode("unicode_escape")


def set_font(run, size: int, bold: bool = False, color=DARK) -> None:
    run.font.name = FONT_NAME
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12.2), Inches(0.8))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = zh(title)
    set_font(r, 24, True, DARK)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = zh(subtitle)
        set_font(r2, 10, False, GRAY)


def add_box(slide, left, top, width, height, title: str, lines: list[str], fill, font_color=WHITE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = zh(title)
    set_font(r, 16, True, font_color)
    for line in lines:
        p = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = zh(line)
        set_font(r, 11, False, font_color)
    return shape


def add_text_card(slide, left, top, width, height, title: str, bullets: list[str]):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT
    shape.line.color.rgb = RGBColor(229, 231, 235)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = zh(title)
    set_font(r, 16, True, DARK)
    for bullet in bullets:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = zh(bullet)
        set_font(r, 11, False, GRAY)
    return shape


def connect(slide, x1, y1, x2, y2, color=GRAY):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(1.5)
    return line


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(
    slide,
    "\\u9879\\u76ee\\u67b6\\u6784\\u56fe",
    "\\u9002\\u7528\\u4e8e\\u5168\\u6808\\u5de5\\u7a0b\\u5e08\\u8003\\u6838\\u5c55\\u793a\\uff1a\\u4ea7\\u54c1\\u3001\\u6280\\u672f\\u3001\\u5de5\\u7a0b\\u5206\\u5c42\\u4e00\\u4f53\\u5316",
)
add_box(slide, Inches(0.5), Inches(1.3), Inches(2.3), Inches(1.1), "\\u7528\\u6237\\u7aef", ["PC Web", "\\u540e\\u7eed\\u53ef\\u6269\\u5c55 App / iPad"], BLUE)
add_box(slide, Inches(3.2), Inches(1.3), Inches(2.5), Inches(1.1), "\\u524d\\u7aef\\u5e94\\u7528\\u5c42", ["React + TypeScript", "Vite + Ant Design"], GREEN)
add_box(slide, Inches(6.15), Inches(1.3), Inches(2.5), Inches(1.1), "\\u540e\\u7aef\\u670d\\u52a1\\u5c42", ["FastAPI", "Router / Service / CRUD"], ORANGE)
add_box(slide, Inches(9.1), Inches(1.3), Inches(2.7), Inches(1.1), "\\u6570\\u636e\\u4e0e\\u57fa\\u7840\\u8bbe\\u65bd", ["MySQL + Redis", "Nginx + Docker\\uff08\\u89c4\\u5212\\uff09"], RED)
connect(slide, Inches(2.8), Inches(1.85), Inches(3.2), Inches(1.85))
connect(slide, Inches(5.7), Inches(1.85), Inches(6.15), Inches(1.85))
connect(slide, Inches(8.65), Inches(1.85), Inches(9.1), Inches(1.85))
add_text_card(slide, Inches(0.7), Inches(3.0), Inches(3.5), Inches(2.5), "\\u6838\\u5fc3\\u4e1a\\u52a1\\u6a21\\u5757", [
    "1. \\u7528\\u6237\\u767b\\u5f55\\u6ce8\\u518c",
    "2. \\u6211\\u7684\\u8863\\u6a71\\uff08\\u8863\\u7269 CRUD\\uff09",
    "3. \\u6211\\u7684\\u7a7f\\u642d\\uff08\\u7ec4\\u5408\\u7ba1\\u7406\\uff09",
    "4. \\u7a7f\\u642d\\u5206\\u4eab\\uff08\\u7ad9\\u5185 + \\u4e00\\u952e\\u5206\\u4eab\\u7ed9\\u95fa\\u871c\\uff09",
    "5. \\u7075\\u611f\\u7a7f\\u642d / \\u5546\\u54c1\\u63a8\\u8350\\u4f4d",
])
add_text_card(slide, Inches(4.55), Inches(3.0), Inches(3.2), Inches(2.5), "\\u540e\\u7aef\\u804c\\u8d23", [
    "1. \\u53c2\\u6570\\u6821\\u9a8c\\u4e0e\\u63a5\\u53e3\\u8fd4\\u56de",
    "2. \\u4e1a\\u52a1\\u903b\\u8f91\\u7f16\\u6392",
    "3. \\u6570\\u636e\\u5e93\\u5b58\\u53d6",
    "4. Redis \\u7f13\\u5b58\\u70ed\\u70b9\\u6570\\u636e",
    "5. \\u540e\\u7eed\\u652f\\u6301 JWT / \\u4e0a\\u4f20 / \\u6d4b\\u8bd5",
])
add_text_card(slide, Inches(8.1), Inches(3.0), Inches(4.5), Inches(2.5), "\\u5de5\\u7a0b\\u76ee\\u6807", [
    "1. \\u524d\\u540e\\u7aef\\u5206\\u79bb",
    "2. \\u5206\\u5c42\\u6e05\\u6670\\uff0c\\u53ef\\u89e3\\u91ca",
    "3. \\u4fbf\\u4e8e\\u6d4b\\u8bd5\\u4e0e\\u6269\\u5c55",
    "4. \\u53ef\\u4f5c\\u4e3a\\u8003\\u6838\\u9879\\u76ee\\u5c55\\u793a",
    "5. \\u540e\\u7eed\\u53ef\\u4ea7\\u54c1\\u5316\\u7ee7\\u7eed\\u6f14\\u8fdb",
])

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "\\u9875\\u9762\\u4e0e\\u529f\\u80fd\\u6e05\\u5355", "\\u6309\\u8003\\u6838\\u7248 MVP \\u548c\\u540e\\u7eed\\u6269\\u5c55\\u62c6\\u5206\\uff0c\\u4fbf\\u4e8e\\u89e3\\u91ca\\u8303\\u56f4\\u8fb9\\u754c")
add_text_card(slide, Inches(0.5), Inches(1.1), Inches(4.0), Inches(5.8), "\\u6838\\u5fc3\\u9875\\u9762", [
    "1. \\u767b\\u5f55\\u9875\\uff1a\\u767b\\u5f55 / \\u6ce8\\u518c / \\u767b\\u5f55\\u6001",
    "2. \\u4eea\\u8868\\u76d8\\uff1a\\u8863\\u7269\\u6570\\u3001\\u7a7f\\u642d\\u6570\\u3001\\u6700\\u8fd1\\u65b0\\u589e",
    "3. \\u6211\\u7684\\u8863\\u6a71\\uff1a\\u5217\\u8868\\u3001\\u7b5b\\u9009\\u3001CRUD",
    "4. \\u6211\\u7684\\u7a7f\\u642d\\uff1a\\u521b\\u5efa\\u7ec4\\u5408\\u3001\\u6807\\u7b7e\\u3001CRUD",
    "5. \\u7a7f\\u642d\\u5206\\u4eab\\uff1a\\u53d1\\u5e03\\u3001\\u516c\\u5f00 / \\u79c1\\u57df\\u5206\\u4eab",
])
add_text_card(slide, Inches(4.7), Inches(1.1), Inches(4.0), Inches(5.8), "\\u4ea7\\u54c1\\u4eae\\u70b9", [
    "1. \\u7075\\u611f\\u7a7f\\u642d\\uff1a\\u65b0\\u7528\\u6237\\u51b7\\u542f\\u52a8\\u5165\\u53e3",
    "2. \\u5546\\u54c1\\u63a8\\u8350\\u4f4d\\uff1a\\u5148\\u642d\\u914d\\uff0c\\u518d\\u8d2d\\u4e70",
    "3. \\u5e7f\\u544a\\u4f4d\\uff1a\\u54c1\\u724c\\u5408\\u4f5c\\u4e0e\\u4e13\\u9898\\u5185\\u5bb9",
    "4. \\u95fa\\u871c\\u5206\\u4eab\\uff1a\\u5206\\u4eab\\u5361\\u7247\\u3001\\u94fe\\u63a5\\u3001\\u5f81\\u6c42\\u610f\\u89c1",
    "5. \\u5185\\u5bb9\\u6c89\\u6dc0\\uff1a\\u4e2a\\u4eba\\u4e3b\\u9875\\u4e0e\\u98ce\\u683c\\u5c55\\u793a",
])
add_text_card(slide, Inches(8.9), Inches(1.1), Inches(3.9), Inches(5.8), "\\u7b2c\\u4e00\\u7248\\u8303\\u56f4\\u63a7\\u5236", [
    "\\u5148\\u505a\\u5b8c\\u6574\\u95ed\\u73af\\uff0c\\u4e0d\\u8ffd\\u6c42\\u529f\\u80fd\\u6700\\u591a",
    "\\u4f18\\u5148\\u5c55\\u793a\\u5de5\\u7a0b\\u89c4\\u8303\\u548c\\u4ea7\\u54c1\\u601d\\u8003",
    "AI \\u751f\\u56fe\\u5148\\u653e\\u89c4\\u5212\\uff0c\\u4e0d\\u505a\\u7b2c\\u4e00\\u6838\\u5fc3",
    "Redis / Docker \\u4f5c\\u4e3a\\u5de5\\u7a0b\\u5316\\u52a0\\u5206\\u9879",
    "\\u5206\\u4eab\\u4e0e\\u5546\\u4e1a\\u5316\\u4f5c\\u4e3a\\u4eae\\u70b9\\uff0c\\u4e0d\\u5f71\\u54cd\\u4e3b\\u94fe\\u8def",
])

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "\\u524d\\u540e\\u7aef\\u5206\\u5c42\\u7ed3\\u6784", "\\u91cd\\u70b9\\u5c55\\u793a\\u4ee3\\u7801\\u7ec4\\u7ec7\\u3001\\u804c\\u8d23\\u8fb9\\u754c\\u548c\\u540e\\u7eed\\u53ef\\u7ef4\\u62a4\\u6027")
add_text_card(slide, Inches(0.5), Inches(1.2), Inches(5.9), Inches(5.6), "\\u524d\\u7aef\\u7ed3\\u6784\\uff08frontend\\uff09", [
    "pages: \\u767b\\u5f55\\u3001\\u8863\\u6a71\\u3001\\u7a7f\\u642d\\u3001\\u5206\\u4eab\\u3001\\u7075\\u611f\\u7a7f\\u642d",
    "components: \\u901a\\u7528\\u7ec4\\u4ef6\\u4e0e\\u4e1a\\u52a1\\u7ec4\\u4ef6",
    "api: \\u8bf7\\u6c42\\u5c01\\u88c5\\uff0c\\u7edf\\u4e00\\u5bf9\\u63a5\\u540e\\u7aef\\u63a5\\u53e3",
    "router: \\u8def\\u7531\\u5b9a\\u4e49\\u4e0e\\u9875\\u9762\\u5165\\u53e3",
    "types: TypeScript \\u7c7b\\u578b\\u5b9a\\u4e49",
    "utils: \\u5de5\\u5177\\u65b9\\u6cd5\\u4e0e\\u516c\\u5171\\u51fd\\u6570",
    "store/hooks: \\u524d\\u7aef\\u72b6\\u6001\\u4e0e\\u590d\\u7528\\u903b\\u8f91",
])
add_text_card(slide, Inches(6.9), Inches(1.2), Inches(5.9), Inches(5.6), "\\u540e\\u7aef\\u7ed3\\u6784\\uff08backend/app\\uff09", [
    "api/routes: \\u5b9a\\u4e49 HTTP \\u8def\\u7531",
    "schemas: Pydantic \\u8bf7\\u6c42\\u4e0e\\u54cd\\u5e94\\u6a21\\u578b",
    "services: \\u4e1a\\u52a1\\u903b\\u8f91\\u7f16\\u6392",
    "crud: \\u6570\\u636e\\u5e93\\u8bfb\\u5199",
    "models: SQLAlchemy \\u8868\\u7ed3\\u6784",
    "db: MySQL / Redis \\u8fde\\u63a5\\u7ba1\\u7406",
    "core: \\u914d\\u7f6e\\u3001\\u5f02\\u5e38\\u3001\\u540e\\u7eed\\u5b89\\u5168\\u80fd\\u529b",
])

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "\\u8fd0\\u884c\\u94fe\\u8def\\u4e0e Redis \\u89d2\\u8272", "\\u8bf4\\u660e\\u4e3a\\u4ec0\\u4e48\\u9879\\u76ee\\u91cc\\u65e2\\u6709 MySQL\\uff0c\\u4e5f\\u8981\\u5f15\\u5165 Redis")
add_box(slide, Inches(0.6), Inches(1.5), Inches(2.1), Inches(0.95), "\\u6d4f\\u89c8\\u5668", ["localhost:5173"], BLUE)
add_box(slide, Inches(3.0), Inches(1.5), Inches(2.2), Inches(0.95), "\\u524d\\u7aef\\u9875\\u9762", ["React \\u9875\\u9762 / \\u8868\\u5355 / \\u5217\\u8868"], GREEN)
add_box(slide, Inches(5.5), Inches(1.5), Inches(2.2), Inches(0.95), "FastAPI \\u63a5\\u53e3", ["REST API"], ORANGE)
add_box(slide, Inches(8.0), Inches(1.2), Inches(2.0), Inches(0.95), "MySQL", ["\\u6b63\\u5f0f\\u4e1a\\u52a1\\u6570\\u636e"], RED)
add_box(slide, Inches(8.0), Inches(2.45), Inches(2.0), Inches(0.95), "Redis", ["\\u5217\\u8868\\u7f13\\u5b58 / \\u4f1a\\u8bdd / \\u70ed\\u70b9"], PURPLE)
add_box(slide, Inches(10.4), Inches(1.85), Inches(2.1), Inches(0.95), "\\u672a\\u6765\\u6269\\u5c55", ["Nginx / Docker / AI"], DARK)
connect(slide, Inches(2.7), Inches(1.98), Inches(3.0), Inches(1.98))
connect(slide, Inches(5.2), Inches(1.98), Inches(5.5), Inches(1.98))
connect(slide, Inches(7.7), Inches(1.98), Inches(8.0), Inches(1.68))
connect(slide, Inches(7.7), Inches(1.98), Inches(8.0), Inches(2.92))
connect(slide, Inches(10.0), Inches(1.68), Inches(10.4), Inches(2.32))
add_text_card(slide, Inches(0.8), Inches(3.4), Inches(3.8), Inches(2.2), "MySQL \\u7684\\u804c\\u8d23", [
    "\\u6b63\\u5f0f\\u5b58\\u50a8\\u7528\\u6237\\u3001\\u8863\\u7269\\u3001\\u7a7f\\u642d\\u3001\\u5206\\u4eab\\u6570\\u636e",
    "\\u4fdd\\u8bc1\\u6570\\u636e\\u5b8c\\u6574\\u6027\\u4e0e\\u5173\\u7cfb\\u7ea6\\u675f",
    "\\u4f5c\\u4e3a\\u6700\\u7ec8\\u53ef\\u4fe1\\u4e1a\\u52a1\\u6570\\u636e\\u6e90",
])
add_text_card(slide, Inches(4.8), Inches(3.4), Inches(3.8), Inches(2.2), "Redis \\u7684\\u804c\\u8d23", [
    "\\u7f13\\u5b58\\u8863\\u7269\\u5217\\u8868\\u7b49\\u9ad8\\u9891\\u8bfb\\u53d6\\u7ed3\\u679c",
    "\\u540e\\u7eed\\u53ef\\u5b58\\u9a8c\\u8bc1\\u7801\\u3001\\u4f1a\\u8bdd\\u3001token \\u9ed1\\u540d\\u5355",
    "\\u540e\\u7eed\\u53ef\\u505a\\u70b9\\u8d5e\\u6570\\u3001\\u6d4f\\u89c8\\u6570\\u3001\\u6392\\u884c\\u699c",
])
add_text_card(slide, Inches(8.8), Inches(3.4), Inches(3.8), Inches(2.2), "\\u5f53\\u524d\\u5df2\\u843d\\u5730", [
    "\\u524d\\u540e\\u7aef\\u672c\\u5730\\u53ef\\u8fd0\\u884c",
    "MySQL \\u5df2\\u63a5\\u5165",
    "Redis \\u5df2\\u9884\\u7559\\u5e76\\u63a5\\u5165\\u7f13\\u5b58\\u5c42",
    "\\u524d\\u7aef\\u4e2d\\u6587\\u9875\\u9762\\u5df2\\u63a5\\u771f\\u5b9e\\u63a5\\u53e3",
])

# Slide 5 ER
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "\\u6570\\u636e\\u5e93 ER \\u56fe\\uff08\\u7b2c\\u4e00\\u7248\\u8bbe\\u8ba1\\uff09", "\\u4ee5\\u7528\\u6237\\u3001\\u8863\\u7269\\u3001\\u7a7f\\u642d\\u3001\\u5206\\u4eab\\u4e3a\\u4e3b\\u7ebf\\u7684\\u6838\\u5fc3\\u5173\\u7cfb")
add_box(slide, Inches(0.4), Inches(1.4), Inches(2.0), Inches(1.7), "users", [
    "id (PK)",
    "username",
    "email",
    "password_hash",
], BLUE)
add_box(slide, Inches(2.9), Inches(1.3), Inches(2.1), Inches(1.9), "clothes", [
    "id (PK)",
    "user_id (FK)",
    "name / category",
    "color / season",
    "image_url",
], GREEN)
add_box(slide, Inches(5.4), Inches(1.3), Inches(2.1), Inches(1.9), "outfits", [
    "id (PK)",
    "user_id (FK)",
    "title / scene",
    "season / cover",
], ORANGE)
add_box(slide, Inches(7.9), Inches(1.3), Inches(2.1), Inches(1.9), "outfit_items", [
    "id (PK)",
    "outfit_id (FK)",
    "clothes_id (FK)",
    "item_type",
], PURPLE)
add_box(slide, Inches(10.4), Inches(1.3), Inches(2.1), Inches(1.9), "posts", [
    "id (PK)",
    "user_id (FK)",
    "outfit_id (FK)",
    "title / content",
    "visibility",
], RED)
add_box(slide, Inches(8.6), Inches(4.0), Inches(1.8), Inches(1.4), "comments", [
    "id (PK)",
    "post_id (FK)",
    "user_id (FK)",
], RGBColor(20, 184, 166))
add_box(slide, Inches(10.8), Inches(4.0), Inches(1.8), Inches(1.4), "post_likes", [
    "id (PK)",
    "post_id (FK)",
    "user_id (FK)",
], RGBColor(236, 72, 153))
connect(slide, Inches(2.4), Inches(2.2), Inches(2.9), Inches(2.2))
connect(slide, Inches(2.4), Inches(2.0), Inches(5.4), Inches(2.0))
connect(slide, Inches(5.0), Inches(2.2), Inches(5.4), Inches(2.2))
connect(slide, Inches(7.5), Inches(2.2), Inches(7.9), Inches(2.2))
connect(slide, Inches(10.0), Inches(2.2), Inches(10.4), Inches(2.2))
connect(slide, Inches(11.4), Inches(3.2), Inches(9.5), Inches(4.0))
connect(slide, Inches(11.4), Inches(3.2), Inches(11.7), Inches(4.0))
add_text_card(slide, Inches(0.6), Inches(4.2), Inches(7.0), Inches(1.8), "\\u5173\\u7cfb\\u8bf4\\u660e", [
    "users \\u4e0e clothes / outfits / posts \\u5747\\u4e3a\\u4e00\\u5bf9\\u591a",
    "outfits \\u4e0e clothes \\u901a\\u8fc7 outfit_items \\u5f62\\u6210\\u591a\\u5bf9\\u591a",
    "posts \\u4e0e comments / post_likes \\u5747\\u4e3a\\u4e00\\u5bf9\\u591a",
])

# Slide 6 speech
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "\\u8003\\u6838\\u7b54\\u8fa9\\u7a3f / \\u8bb2\\u89e3\\u8bdd\\u672f", "\\u5efa\\u8bae\\u5148\\u8bb2\\u4e1a\\u52a1\\uff0c\\u518d\\u8bb2\\u67b6\\u6784\\uff0c\\u6700\\u540e\\u8bb2\\u5de5\\u7a0b\\u4e0e\\u6269\\u5c55")
add_text_card(slide, Inches(0.5), Inches(1.15), Inches(4.0), Inches(5.7), "\\u5f00\\u573a\\u8bdd\\u672f", [
    "\\u8fd9\\u662f\\u4e00\\u4e2a\\u9762\\u5411\\u4e2a\\u4eba\\u7528\\u6237\\u7684\\u7535\\u5b50\\u8863\\u6a71\\u4e0e\\u7a7f\\u642d\\u7ba1\\u7406\\u7cfb\\u7edf\\u3002",
    "\\u5b83\\u7684\\u4e3b\\u7ebf\\u662f\\u5e2e\\u52a9\\u7528\\u6237\\u7ba1\\u7406\\u8863\\u7269\\u3001\\u7ec4\\u5408\\u7a7f\\u642d\\u3001\\u53d1\\u5e03\\u5206\\u4eab\\uff0c",
    "\\u540c\\u65f6\\u4fdd\\u7559\\u4e86\\u7075\\u611f\\u7a7f\\u642d\\u3001\\u5546\\u54c1\\u63a8\\u8350\\u548c AI \\u80fd\\u529b\\u7684\\u6269\\u5c55\\u7a7a\\u95f4\\u3002",
])
add_text_card(slide, Inches(4.7), Inches(1.15), Inches(4.0), Inches(5.7), "\\u6280\\u672f\\u8bdd\\u672f", [
    "\\u524d\\u7aef\\u91c7\\u7528 React + TypeScript + Vite + Ant Design\\uff0c",
    "\\u540e\\u7aef\\u91c7\\u7528 FastAPI + SQLAlchemy + MySQL\\uff0c",
    "\\u6574\\u4f53\\u662f\\u524d\\u540e\\u7aef\\u5206\\u79bb\\u67b6\\u6784\\u3002",
    "\\u540e\\u7aef\\u4f7f\\u7528 Route / Schema / Service / CRUD \\u5206\\u5c42\\uff0c\\u4fbf\\u4e8e\\u7ef4\\u62a4\\u548c\\u6d4b\\u8bd5\\u3002",
])
add_text_card(slide, Inches(8.9), Inches(1.15), Inches(3.9), Inches(5.7), "\\u5de5\\u7a0b\\u4e0e\\u6269\\u5c55\\u8bdd\\u672f", [
    "MySQL \\u8d1f\\u8d23\\u6b63\\u5f0f\\u4e1a\\u52a1\\u6570\\u636e\\uff0cRedis \\u8d1f\\u8d23\\u70ed\\u70b9\\u7f13\\u5b58\\u3002",
    "\\u5f53\\u524d\\u9879\\u76ee\\u5df2\\u5177\\u5907\\u672c\\u5730\\u8fd0\\u884c\\u80fd\\u529b\\uff0c",
    "\\u540e\\u7eed\\u53ef\\u4ee5\\u7ee7\\u7eed\\u63a5\\u5165 JWT\\u3001Docker\\u3001Nginx \\u548c AI \\u6a21\\u5757\\u3002",
    "\\u8fd9\\u4e2a\\u9879\\u76ee\\u4e0d\\u53ea\\u662f demo\\uff0c\\u800c\\u662f\\u53ef\\u4ee5\\u7ee7\\u7eed\\u4ea7\\u54c1\\u5316\\u6f14\\u8fdb\\u7684\\u57fa\\u7840\\u7248\\u672c\\u3002",
])

output = r"D:\\Project\\dress\\docs\\dress-architecture-final.pptx"
prs.save(output)
print(output)
